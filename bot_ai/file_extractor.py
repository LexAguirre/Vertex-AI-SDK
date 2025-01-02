from pathlib import Path

import pandas as pd
from django.conf import settings
from django.utils import timezone
from docx import Document
from langchain_community.document_loaders import PyPDFLoader
from pptx import Presentation

from app.bot_ai.utils import get_file_divition


class PDFExtractor:
    """
    A class that handles the extraction of text and the management of files from different formats (PDF, PPTX, DOCX, XLSX).
    It provides utilities for converting, extracting, and exporting files for a specific customer.

    Attributes:
        file_text (str): A variable to store the extracted text.
        customer_name (str): The name of the customer for whom files are being processed.
        customer_id (int): The unique ID of the customer.
    """  # noqa: E501

    def __init__(self, customer_name, customer_id):
        """
        Initializes the PDFExtractor with customer details.

        Args:
            customer_name (str): The name of the customer.
            customer_id (int): The ID of the customer.
        """
        self.file_text = ""
        self.customer_name = customer_name
        self.customer_id = customer_id

    def file_name(self):
        """
        Generates a filename based on the current time and customer name.

        Returns:
            str: A formatted filename containing the customer's name and the current timestamp.
        """  # noqa: E501
        time_name = timezone.now().strftime("%Y-%m-%d--%H-%M-%S")
        return f"{self.customer_name}-date-{time_name}"

    def customer_folder(self):
        """
        Creates a folder name for the customer based on their name and ID.

        Returns:
            str: A formatted folder name.
        """
        return f"{self.customer_name}_ID_{self.customer_id}"

    def importation_manager(self):
        """
        Creates the necessary import folder structure for a customer within the media directory.
        """  # noqa: E501
        customer_folder = self.customer_folder()
        import_folder = Path(settings.MEDIA_ROOT) / "import" / customer_folder
        import_folder.mkdir(parents=True, exist_ok=True)

    def create_export_folder(self):
        """
        Creates the necessary export folder structure for a customer within the media directory.
        """  # noqa: E501
        customer_folder = self.customer_folder()
        export_folder = Path(settings.MEDIA_ROOT) / "export" / customer_folder
        export_folder.mkdir(parents=True, exist_ok=True)

    def extract_text_pdf(self):
        """
        Extracts text from a PDF file and stores it in the `file_text` attribute.
        """
        self.create_export_folder()
        customer_folder = self.customer_folder()
        scrape_file = "Tupper_Tips_NORTE.pdf"
        url_doc = f"app/media/import/{customer_folder}/{scrape_file}"
        loader = PyPDFLoader(url_doc)
        pages = loader.load()

        for i in range(len(pages)):
            page = pages[i]
            self.file_text += page.page_content + "\n"

    def extract_text_pptx(self):
        """
        Extracts text from a PowerPoint (.pptx) file and stores it in the `file_text` attribute.
        """  # noqa: E501
        self.create_export_folder()
        customer_folder = self.customer_folder()
        scrape_file = "Seaborn-Scikitlearn.pptx"
        url_doc = f"app/media/import/{customer_folder}/{scrape_file}"

        presentation = Presentation(url_doc)
        self.file_text = ""

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    self.file_text += shape.text + "\n"

    def extract_text_docx(self):
        """
        Extracts text from a Word document (.docx) and stores it in the `file_text` attribute.
        """  # noqa: E501
        self.create_export_folder()
        customer_folder = self.customer_folder()
        scrape_file = "MEMORIA-DE-PP-2024.docx"
        url_doc = f"app/media/import/{customer_folder}/{scrape_file}"

        document = Document(url_doc)
        self.file_text = ""

        for para in document.paragraphs:
            self.file_text += para.text + "\n"

    def convert_xlsx_to_csv(self):
        """
        Converts an Excel (.xlsx) file to a CSV file and saves it in the export folder.
        """
        self.create_export_folder()
        customer_folder = self.customer_folder()
        scrape_file = "Lista de asistencia del club.xlsx"
        url_doc = f"app/media/import/{customer_folder}/{scrape_file}"

        dataframe_xlsx = pd.read_excel(url_doc, engine="openpyxl")
        file_name = self.file_name()
        txt_file_path = f"app/media/export/{customer_folder}/{file_name}.csv"
        dataframe_xlsx.to_csv(txt_file_path, index=False)

    def split_csv(self, input_file):
        """
        Splits a large CSV file into smaller parts and saves each part in the export folder.

        Args:
            input_file (str): The name of the input CSV file.

        Returns:
            list: A list of file paths for the split CSV files.
            list: A list of names for the split CSV files.
        """  # noqa: E501
        self.create_export_folder()
        files_list = []
        files_name_list = []
        customer_folder = self.customer_folder()
        url_doc_import = f"app/media/import/{customer_folder}/{input_file}.csv"
        num_files = get_file_divition(url_doc_import)

        if num_files == 1:
            return [url_doc_import], [f"{input_file}.csv"]

        csv_dataframe = pd.read_csv(url_doc_import)
        total_rows = len(csv_dataframe)

        if total_rows % num_files != 0:
            lines_per_file = (total_rows // num_files) + 1
        else:
            lines_per_file = total_rows // num_files

        for i in range(num_files):
            start_row = i * lines_per_file
            end_row = total_rows if i == num_files - 1 else start_row + lines_per_file
            df_subset = csv_dataframe.iloc[start_row:end_row]
            output_file_name = f"{input_file}_part_{i + 1}"
            url_doc_export = (
                f"app/media/export/{customer_folder}/{output_file_name}.csv"
            )
            df_subset.to_csv(url_doc_export, index=False)
            files_name_list.append(output_file_name)
            files_list.append(url_doc_export)

        return files_list, files_name_list

    def txt_file_manager(self):
        """
        Saves the extracted text (from PDF, PPTX, DOCX) to a text file in the export folder.

        Returns:
            str: The path to the saved text file.
        """  # noqa: E501
        customer_folder = self.customer_folder()
        file_name = self.file_name()
        txt_file_path = f"app/media/export/{customer_folder}/{file_name}.txt"

        with open(txt_file_path, "wt", encoding="utf-8") as file:  # noqa: PTH123, UP015
            file.write(self.file_text)

        return txt_file_path
